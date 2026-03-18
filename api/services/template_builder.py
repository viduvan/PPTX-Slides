"""
Template Builder — Logic for creating slide templates with specific themes and layouts.
Developed by ChimSe (viduvan) - https://github.com/viduvan

Features:
  - Multiple theme presets (dark purple, ocean, forest, sunset, midnight, crimson)
  - Gradient backgrounds per theme
  - Styled typography (Calibri Light titles, Calibri body)
  - Decorative accent bar at slide bottom
  - Auto-fit font sizing so content never overflows
  - Image placement support
"""
import logging
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE

logger = logging.getLogger("odin_api.services.template_builder")

# ── Theme Categories & Registry ─────────────────────────────
THEME_CATEGORIES = {
    "business":    {"label": "Business",    "label_vi": "Doanh nghiệp",    "emoji": "🏢", "order": 1},
    "creative":    {"label": "Creative",    "label_vi": "Sáng tạo",        "emoji": "🎨", "order": 2},
    "education":   {"label": "Education",   "label_vi": "Giáo dục",        "emoji": "📚", "order": 3},
    "technology":  {"label": "Technology",  "label_vi": "Công nghệ",       "emoji": "💻", "order": 4},
}

# ── All 30 Theme Presets ────────────────────────────────────
THEMES = {
    # ═══ BUSINESS (5) ═══════════════════════════════════════
    "corporate_blue": {
        "bg_dark":       RGBColor(0x0A, 0x14, 0x28),
        "bg_gradient":   RGBColor(0x10, 0x20, 0x3C),
        "title":         RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle":      RGBColor(0x93, 0xC5, 0xFD),
        "body":          RGBColor(0xD4, 0xDE, 0xEC),
        "accent":        RGBColor(0x1D, 0x4E, 0xD8),
        "accent_light":  RGBColor(0x3B, 0x82, 0xF6),
        "muted":         RGBColor(0x88, 0x99, 0xAA),
    },
    "executive_gray": {
        "bg_dark":       RGBColor(0x14, 0x14, 0x16),
        "bg_gradient":   RGBColor(0x1E, 0x1E, 0x22),
        "title":         RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle":      RGBColor(0xD4, 0xAF, 0x37),
        "body":          RGBColor(0xE0, 0xE0, 0xE0),
        "accent":        RGBColor(0xB8, 0x96, 0x0C),
        "accent_light":  RGBColor(0xD4, 0xAF, 0x37),
        "muted":         RGBColor(0x9C, 0x9C, 0xA0),
    },
    "finance_green": {
        "bg_dark":       RGBColor(0x06, 0x12, 0x10),
        "bg_gradient":   RGBColor(0x0A, 0x20, 0x1A),
        "title":         RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle":      RGBColor(0x6E, 0xE7, 0xB7),
        "body":          RGBColor(0xD1, 0xFA, 0xE5),
        "accent":        RGBColor(0x05, 0x96, 0x69),
        "accent_light":  RGBColor(0x10, 0xB9, 0x81),
        "muted":         RGBColor(0x7C, 0xA8, 0x96),
    },
    "legal_navy": {
        "bg_dark":       RGBColor(0x08, 0x0C, 0x1A),
        "bg_gradient":   RGBColor(0x0E, 0x14, 0x2C),
        "title":         RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle":      RGBColor(0xA5, 0xB4, 0xFC),
        "body":          RGBColor(0xC7, 0xD2, 0xFE),
        "accent":        RGBColor(0x43, 0x38, 0xCA),
        "accent_light":  RGBColor(0x63, 0x66, 0xF1),
        "muted":         RGBColor(0x8B, 0x8E, 0xB0),
    },
    "consulting_teal": {
        "bg_dark":       RGBColor(0x08, 0x14, 0x18),
        "bg_gradient":   RGBColor(0x0C, 0x20, 0x28),
        "title":         RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle":      RGBColor(0x5E, 0xEA, 0xD4),
        "body":          RGBColor(0xCC, 0xFB, 0xF1),
        "accent":        RGBColor(0x0D, 0x94, 0x88),
        "accent_light":  RGBColor(0x14, 0xB8, 0xA6),
        "muted":         RGBColor(0x80, 0xA8, 0xA0),
    },

    # ═══ CREATIVE (5) ═══════════════════════════════════════
    "bold_orange": {
        "bg_dark":       RGBColor(0x1A, 0x0A, 0x05),
        "bg_gradient":   RGBColor(0x30, 0x10, 0x08),
        "title":         RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle":      RGBColor(0xFD, 0xBA, 0x74),
        "body":          RGBColor(0xF0, 0xE0, 0xD0),
        "accent":        RGBColor(0xEA, 0x58, 0x0C),
        "accent_light":  RGBColor(0xFB, 0x92, 0x3C),
        "muted":         RGBColor(0xBB, 0xA0, 0x8C),
    },
    "artistic_purple": {
        "bg_dark":       RGBColor(0x14, 0x06, 0x1A),
        "bg_gradient":   RGBColor(0x24, 0x08, 0x30),
        "title":         RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle":      RGBColor(0xE8, 0x79, 0xF9),
        "body":          RGBColor(0xF0, 0xD8, 0xF5),
        "accent":        RGBColor(0xA8, 0x55, 0xF7),
        "accent_light":  RGBColor(0xC0, 0x84, 0xFC),
        "muted":         RGBColor(0xA0, 0x88, 0xAA),
    },
    "neon_pop": {
        "bg_dark":       RGBColor(0x0A, 0x0A, 0x0F),
        "bg_gradient":   RGBColor(0x12, 0x12, 0x1A),
        "title":         RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle":      RGBColor(0x22, 0xD3, 0xEE),
        "body":          RGBColor(0xE0, 0xF2, 0xFE),
        "accent":        RGBColor(0xE8, 0x36, 0x8F),
        "accent_light":  RGBColor(0x22, 0xD3, 0xEE),
        "muted":         RGBColor(0x88, 0x99, 0xAA),
    },
    "retro_vintage": {
        "bg_dark":       RGBColor(0x1C, 0x14, 0x0A),
        "bg_gradient":   RGBColor(0x2A, 0x20, 0x12),
        "title":         RGBColor(0xFF, 0xF5, 0xE0),
        "subtitle":      RGBColor(0xF0, 0xC0, 0x70),
        "body":          RGBColor(0xE8, 0xDD, 0xCC),
        "accent":        RGBColor(0xC2, 0x71, 0x0E),
        "accent_light":  RGBColor(0xD9, 0x97, 0x06),
        "muted":         RGBColor(0xA0, 0x94, 0x80),
    },
    "rose_pink": {
        "bg_dark":       RGBColor(0x18, 0x08, 0x14),
        "bg_gradient":   RGBColor(0x28, 0x0C, 0x22),
        "title":         RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle":      RGBColor(0xF9, 0xA8, 0xD4),
        "body":          RGBColor(0xF0, 0xDC, 0xE8),
        "accent":        RGBColor(0xDB, 0x27, 0x77),
        "accent_light":  RGBColor(0xF4, 0x72, 0xB6),
        "muted":         RGBColor(0xB0, 0x8C, 0xA0),
    },

    # ═══ EDUCATION (5) ══════════════════════════════════════
    "scholar_blue": {
        "bg_dark":       RGBColor(0x0C, 0x16, 0x2C),
        "bg_gradient":   RGBColor(0x14, 0x22, 0x40),
        "title":         RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle":      RGBColor(0x7D, 0xD3, 0xFC),
        "body":          RGBColor(0xD6, 0xED, 0xF5),
        "accent":        RGBColor(0x06, 0x8F, 0xCF),
        "accent_light":  RGBColor(0x38, 0xBD, 0xF8),
        "muted":         RGBColor(0x8E, 0xA8, 0xBB),
    },
    "campus_green": {
        "bg_dark":       RGBColor(0x07, 0x15, 0x0B),
        "bg_gradient":   RGBColor(0x0A, 0x28, 0x14),
        "title":         RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle":      RGBColor(0x86, 0xEF, 0xAC),
        "body":          RGBColor(0xD8, 0xF0, 0xDB),
        "accent":        RGBColor(0x16, 0xA3, 0x4A),
        "accent_light":  RGBColor(0x4A, 0xDE, 0x80),
        "muted":         RGBColor(0x8C, 0xAF, 0x94),
    },
    "library_brown": {
        "bg_dark":       RGBColor(0x18, 0x10, 0x08),
        "bg_gradient":   RGBColor(0x28, 0x1C, 0x10),
        "title":         RGBColor(0xFF, 0xF5, 0xE8),
        "subtitle":      RGBColor(0xD4, 0xA5, 0x74),
        "body":          RGBColor(0xE8, 0xDC, 0xCC),
        "accent":        RGBColor(0x92, 0x55, 0x1E),
        "accent_light":  RGBColor(0xB4, 0x70, 0x2E),
        "muted":         RGBColor(0xA0, 0x90, 0x78),
    },
    "science_teal": {
        "bg_dark":       RGBColor(0x04, 0x14, 0x18),
        "bg_gradient":   RGBColor(0x08, 0x22, 0x2A),
        "title":         RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle":      RGBColor(0x67, 0xE8, 0xF9),
        "body":          RGBColor(0xCC, 0xFB, 0xF1),
        "accent":        RGBColor(0x06, 0xB6, 0xD4),
        "accent_light":  RGBColor(0x22, 0xD3, 0xEE),
        "muted":         RGBColor(0x7C, 0xA8, 0xB0),
    },
    "chalkboard": {
        "bg_dark":       RGBColor(0x0A, 0x1A, 0x10),
        "bg_gradient":   RGBColor(0x14, 0x2B, 0x1C),
        "title":         RGBColor(0xFF, 0xFF, 0xF0),
        "subtitle":      RGBColor(0xE0, 0xE0, 0xC8),
        "body":          RGBColor(0xD0, 0xD0, 0xB8),
        "accent":        RGBColor(0xF5, 0xC5, 0x42),
        "accent_light":  RGBColor(0xFA, 0xDB, 0x6A),
        "muted":         RGBColor(0x90, 0xA0, 0x88),
    },


    # ═══ TECHNOLOGY (5) ═════════════════════════════════════
    "cyber_punk": {
        "bg_dark":       RGBColor(0x0A, 0x06, 0x14),
        "bg_gradient":   RGBColor(0x14, 0x08, 0x22),
        "title":         RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle":      RGBColor(0x22, 0xD3, 0xEE),
        "body":          RGBColor(0xE0, 0xE8, 0xF0),
        "accent":        RGBColor(0xE8, 0x36, 0x8F),
        "accent_light":  RGBColor(0x22, 0xD3, 0xEE),
        "muted":         RGBColor(0x88, 0x88, 0xAA),
    },
    "matrix_green": {
        "bg_dark":       RGBColor(0x04, 0x0A, 0x04),
        "bg_gradient":   RGBColor(0x08, 0x14, 0x08),
        "title":         RGBColor(0x00, 0xFF, 0x41),
        "subtitle":      RGBColor(0x66, 0xFF, 0x88),
        "body":          RGBColor(0xBB, 0xEE, 0xCC),
        "accent":        RGBColor(0x00, 0xCC, 0x33),
        "accent_light":  RGBColor(0x00, 0xFF, 0x41),
        "muted":         RGBColor(0x55, 0x99, 0x66),
    },
    "ai_blue": {
        "bg_dark":       RGBColor(0x04, 0x08, 0x1A),
        "bg_gradient":   RGBColor(0x08, 0x10, 0x30),
        "title":         RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle":      RGBColor(0x60, 0xA5, 0xFA),
        "body":          RGBColor(0xBF, 0xDB, 0xFE),
        "accent":        RGBColor(0x25, 0x63, 0xEB),
        "accent_light":  RGBColor(0x38, 0x82, 0xF6),
        "muted":         RGBColor(0x70, 0x88, 0xB0),
    },
    "quantum_violet": {
        "bg_dark":       RGBColor(0x0C, 0x06, 0x18),
        "bg_gradient":   RGBColor(0x18, 0x0C, 0x2E),
        "title":         RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle":      RGBColor(0xA7, 0x8B, 0xFA),
        "body":          RGBColor(0xDD, 0xD6, 0xFE),
        "accent":        RGBColor(0x7C, 0x3A, 0xED),
        "accent_light":  RGBColor(0x0E, 0xA5, 0xC9),
        "muted":         RGBColor(0x8A, 0x80, 0xAA),
    },
    "data_orange": {
        "bg_dark":       RGBColor(0x12, 0x0C, 0x08),
        "bg_gradient":   RGBColor(0x1E, 0x14, 0x0C),
        "title":         RGBColor(0xFF, 0xFF, 0xFF),
        "subtitle":      RGBColor(0xFB, 0xBF, 0x24),
        "body":          RGBColor(0xFE, 0xF3, 0xC7),
        "accent":        RGBColor(0xF5, 0x9E, 0x0B),
        "accent_light":  RGBColor(0xFB, 0xBF, 0x24),
        "muted":         RGBColor(0xA0, 0x90, 0x70),
    },
}

# ── Theme → Category mapping ────────────────────────────────
THEME_REGISTRY = {
    # Business
    "corporate_blue":  {"category": "business",   "label": "Corporate Blue",  "label_vi": "Xanh doanh nghiệp",  "emoji": "💼"},
    "executive_gray":  {"category": "business",   "label": "Executive Gray",  "label_vi": "Xám sang trọng",     "emoji": "🏛️"},
    "finance_green":   {"category": "business",   "label": "Finance Green",   "label_vi": "Xanh tài chính",     "emoji": "📊"},
    "legal_navy":      {"category": "business",   "label": "Legal Navy",      "label_vi": "Xanh hải quân",      "emoji": "⚖️"},
    "consulting_teal": {"category": "business",   "label": "Consulting Teal", "label_vi": "Xanh ngọc tư vấn",   "emoji": "🤝"},
    # Creative
    "bold_orange":     {"category": "creative",   "label": "Bold Orange",     "label_vi": "Cam nổi bật",        "emoji": "🔥"},
    "artistic_purple": {"category": "creative",   "label": "Artistic Purple", "label_vi": "Tím nghệ thuật",     "emoji": "🎭"},
    "neon_pop":        {"category": "creative",   "label": "Neon Pop",        "label_vi": "Neon rực rỡ",        "emoji": "⚡"},
    "retro_vintage":   {"category": "creative",   "label": "Retro Vintage",   "label_vi": "Cổ điển hoài niệm",  "emoji": "📻"},
    "rose_pink":       {"category": "creative",   "label": "Rose Pink",       "label_vi": "Hồng rose",          "emoji": "🌸"},
    # Education
    "scholar_blue":    {"category": "education",  "label": "Scholar Blue",    "label_vi": "Xanh học thuật",     "emoji": "🎓"},
    "campus_green":    {"category": "education",  "label": "Campus Green",    "label_vi": "Xanh campus",        "emoji": "🌿"},
    "library_brown":   {"category": "education",  "label": "Library Brown",   "label_vi": "Nâu thư viện",       "emoji": "📖"},
    "science_teal":    {"category": "education",  "label": "Science Teal",    "label_vi": "Xanh khoa học",      "emoji": "🔬"},
    "chalkboard":      {"category": "education",  "label": "Chalkboard",      "label_vi": "Bảng phấn",          "emoji": "📝"},
    # Technology
    "cyber_punk":      {"category": "technology", "label": "Cyber Punk",      "label_vi": "Cyber Punk",         "emoji": "🤖"},
    "matrix_green":    {"category": "technology", "label": "Matrix Green",    "label_vi": "Xanh Matrix",        "emoji": "💚"},
    "ai_blue":         {"category": "technology", "label": "AI Blue",         "label_vi": "Xanh AI",            "emoji": "🧠"},
    "quantum_violet":  {"category": "technology", "label": "Quantum Violet",  "label_vi": "Tím lượng tử",       "emoji": "🌌"},
    "data_orange":     {"category": "technology", "label": "Data Orange",     "label_vi": "Cam dữ liệu",       "emoji": "📡"},
}

AVAILABLE_THEMES = list(THEMES.keys())
DEFAULT_THEME = "corporate_blue"

# ── Slide Dimensions (16:9) ─────────────────────────────────
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ── Font Configuration ──────────────────────────────────────
FONT_TITLE = "Calibri Light"
FONT_BODY = "Calibri"

TITLE_SLIDE_TITLE_SIZE = Pt(40)
TITLE_SLIDE_SUBTITLE_SIZE = Pt(20)

CONTENT_TITLE_SIZE = Pt(28)
CONTENT_BODY_SIZE = Pt(16)
CONTENT_BODY_MIN_SIZE = Pt(10)

ACCENT_BAR_HEIGHT = Inches(0.08)


def get_theme(theme_name: str | None = None) -> dict:
    """Get a theme palette by name. Falls back to default if not found."""
    if theme_name and theme_name in THEMES:
        return THEMES[theme_name]
    # Try fuzzy match
    if theme_name:
        name = theme_name.lower().strip().replace(" ", "_").replace("-", "_")
        if name in THEMES:
            return THEMES[name]
        # Partial match
        for key in THEMES:
            if name in key or key in name:
                return THEMES[key]
    return THEMES[DEFAULT_THEME]


def _set_slide_gradient(slide, colors):
    """Apply dark gradient background to a slide."""
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_angle = 315

    stop0 = fill.gradient_stops[0]
    stop0.position = 0.0
    stop0.color.rgb = colors["bg_dark"]

    stop1 = fill.gradient_stops[1]
    stop1.position = 1.0
    stop1.color.rgb = colors["bg_gradient"]


def _add_accent_bar(slide, colors, y_position=None):
    """Add decorative accent gradient bar at the bottom."""
    if y_position is None:
        y_position = SLIDE_HEIGHT - ACCENT_BAR_HEIGHT - Inches(0.3)

    bar = slide.shapes.add_shape(
        1, Inches(0.5), y_position,
        SLIDE_WIDTH - Inches(1.0), ACCENT_BAR_HEIGHT,
    )
    bar.line.fill.background()

    fill = bar.fill
    fill.gradient()
    fill.gradient_angle = 0
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[0].color.rgb = colors["accent"]
    fill.gradient_stops[1].position = 1.0
    fill.gradient_stops[1].color.rgb = colors["accent_light"]


def _add_decorative_shape(slide, x, y, w, h, color, shape_type=1):
    """Add a small solid-color decorative shape (rectangle or oval)."""
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _estimate_text_lines(text, chars_per_line=70):
    """Estimate line count for text block."""
    total = 0
    for line in text.split("\n"):
        if not line.strip():
            total += 1
        else:
            total += max(1, len(line) // chars_per_line + 1)
    return total


def _calculate_font_size(content_text, available_height_inches,
                         max_size=CONTENT_BODY_SIZE,
                         min_size=CONTENT_BODY_MIN_SIZE):
    """Auto-shrink font size so content fits within available height."""
    LINE_HEIGHT_RATIO = 0.028  # inches per pt per line

    current_pt = max_size.pt if hasattr(max_size, "pt") else max_size
    min_pt = min_size.pt if hasattr(min_size, "pt") else min_size

    while current_pt >= min_pt:
        chars_per_line = int(80 * (16 / current_pt))
        estimated_lines = _estimate_text_lines(content_text, chars_per_line)
        total_height = estimated_lines * current_pt * LINE_HEIGHT_RATIO

        if total_height <= available_height_inches:
            return Pt(current_pt)
        current_pt -= 1

    return Pt(min_pt)


def _strip_html_and_markdown(text):
    """Remove HTML tags and markdown formatting from text."""
    text = re.sub(r'<[^>]+>', '', text)
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)
    text = re.sub(r'^#{1,6}\s+', '', text, flags=re.MULTILINE)
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


def _format_content_text(text_frame, content, font_size, colors):
    """Format content with bullet points and styled paragraphs."""
    text_frame.clear()
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.word_wrap = True

    content = _strip_html_and_markdown(content)

    lines = content.split("\n")
    first = True

    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue

        para = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
        first = False

        is_bullet = stripped.startswith(("- ", "• ", "* ", "– "))
        if is_bullet:
            stripped = stripped.lstrip("-•*– ").strip()
            run = para.add_run()
            run.text = f"  •  {stripped}"
            para.space_before = Pt(2)
            para.space_after = Pt(2)
        else:
            run = para.add_run()
            run.text = stripped
            para.space_before = Pt(4)
            para.space_after = Pt(4)

        run.font.name = FONT_BODY
        run.font.size = font_size
        run.font.color.rgb = colors["body"]


def build_title_slide(prs, title, subtitle="", colors=None):
    """Create a visually striking title slide."""
    if colors is None:
        colors = THEMES[DEFAULT_THEME]

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    _set_slide_gradient(slide, colors)

    # Small accent decoration top-left
    _add_decorative_shape(slide, Inches(0.3), Inches(0.3),
                          Inches(0.5), Pt(4), colors["accent"])

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(2.0), Inches(11.3), Inches(2.0)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.name = FONT_TITLE
    run.font.size = TITLE_SLIDE_TITLE_SIZE
    run.font.color.rgb = colors["title"]
    run.font.bold = True

    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(
            Inches(2.0), Inches(4.2), Inches(9.3), Inches(1.0)
        )
        tf2 = sub_box.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.name = FONT_BODY
        run2.font.size = TITLE_SLIDE_SUBTITLE_SIZE
        run2.font.color.rgb = colors["subtitle"]

    _add_accent_bar(slide, colors, Inches(5.5))
    return slide


def build_content_slide(prs, title, content, image_path=None,
                        slide_number=None, colors=None):
    """
    Create a content slide with title, body, optional image.
    Auto-adjusts font size to fit content.
    """
    if colors is None:
        colors = THEMES[DEFAULT_THEME]

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_gradient(slide, colors)

    # Slide number badge (top-right)
    if slide_number is not None:
        num_box = slide.shapes.add_textbox(
            SLIDE_WIDTH - Inches(1.2), Inches(0.3), Inches(0.8), Inches(0.4)
        )
        p_num = num_box.text_frame.paragraphs[0]
        p_num.alignment = PP_ALIGN.RIGHT
        run_num = p_num.add_run()
        run_num.text = str(slide_number)
        run_num.font.name = FONT_BODY
        run_num.font.size = Pt(12)
        run_num.font.color.rgb = colors["muted"]

    # If image, split layout 60/40
    has_image = image_path and Path(image_path).exists()
    if has_image:
        content_width = Inches(7.5)
        content_left = Inches(0.8)
        img_left = Inches(8.8)
        img_width = Inches(4.0)
        img_top = Inches(1.8)
        img_height = Inches(4.5)
    else:
        content_width = Inches(11.5)
        content_left = Inches(0.8)

    # Title
    title_box = slide.shapes.add_textbox(
        content_left, Inches(0.5), content_width, Inches(1.0)
    )
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    p_title = tf_title.paragraphs[0]
    p_title.alignment = PP_ALIGN.LEFT
    run_title = p_title.add_run()
    run_title.text = title
    run_title.font.name = FONT_TITLE
    run_title.font.size = CONTENT_TITLE_SIZE
    run_title.font.color.rgb = colors["title"]
    run_title.font.bold = True

    # Title underline accent
    underline = slide.shapes.add_shape(
        1, content_left, Inches(1.45), Inches(2.0), Pt(3),
    )
    underline.line.fill.background()
    underline.fill.solid()
    underline.fill.fore_color.rgb = colors["accent"]

    # Body content
    body_top = Inches(1.7)
    body_height = 5.0
    body_box = slide.shapes.add_textbox(
        content_left, body_top, content_width, Inches(body_height)
    )
    tf_body = body_box.text_frame
    tf_body.word_wrap = True

    font_size = _calculate_font_size(content, body_height)
    logger.debug(f"Slide '{title}': auto-fit font size = {font_size.pt}pt")
    _format_content_text(tf_body, content, font_size, colors)

    # Add image if available
    if has_image:
        try:
            slide.shapes.add_picture(
                str(image_path), img_left, img_top, img_width, img_height
            )
            logger.info(f"Added image to slide: {image_path}")
        except Exception as e:
            logger.warning(f"Failed to add image: {e}")

    _add_accent_bar(slide, colors)
    return slide


def build_themed_presentation(slides_data=None, image_paths=None,
                              theme_name=None):
    """
    Build a complete themed presentation.

    Args:
        slides_data: List of dicts with: slide_number, title, content, narration
        image_paths: Dict mapping slide_number -> image file path
        theme_name: Theme preset name (e.g. 'ocean', 'forest', 'sunset')

    Returns:
        Presentation object
    """
    colors = get_theme(theme_name)
    used_theme = theme_name or DEFAULT_THEME
    logger.info(f"Building presentation with theme: {used_theme}")

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    if not slides_data:
        return prs

    if image_paths is None:
        image_paths = {}

    # First slide as title slide
    first = slides_data[0]
    subtitle = first.get("content", "")
    if len(subtitle) > 120:
        subtitle = subtitle[:120] + "..."
    build_title_slide(prs, title=first.get("title", "Presentation"),
                      subtitle=subtitle, colors=colors)

    # Remaining slides as content
    for i, sd in enumerate(slides_data[1:], start=2):
        img = image_paths.get(sd.get("slide_number"), None)
        build_content_slide(
            prs,
            title=sd.get("title", f"Slide {i}"),
            content=sd.get("content", ""),
            image_path=img,
            slide_number=i,
            colors=colors,
        )

    return prs
