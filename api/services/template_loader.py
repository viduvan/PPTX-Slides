"""
Template Loader — Smart loader that auto-detects slide layouts from .pptx template files.
Caches loaded templates in memory for performance.

The loader:
1. Opens a .pptx template file
2. Auto-detects which slide_layout is best for "Title" vs "Content"
3. Creates new slides by cloning layouts and populating text placeholders
4. Caches parsed template metadata to avoid re-reading large files
"""
from __future__ import annotations

import copy
import io
import logging
from pathlib import Path
from functools import lru_cache

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .template_builder import (
    THEMES, THEME_REGISTRY, THEME_CATEGORIES,
    AVAILABLE_THEMES, DEFAULT_THEME,
)
from ..core.config import settings

logger = logging.getLogger("pptx_api.template_loader")

# ── In-memory cache ─────────────────────────────────────────
# Cache: theme_id -> bytes (raw .pptx file read into memory)
_template_bytes_cache: dict[str, bytes] = {}


def _resolve_template_path(theme_id: str) -> Path | None:
    """Find the .pptx file for a theme_id in the templates directory."""
    reg = THEME_REGISTRY.get(theme_id)
    if not reg:
        return None

    category = reg["category"]
    path = settings.TEMPLATES_DIR / category / f"{theme_id}.pptx"
    if path.exists():
        return path

    # Fallback: search anywhere in templates/
    for f in settings.TEMPLATES_DIR.rglob(f"{theme_id}.pptx"):
        return f

    return None


def _load_template_bytes(theme_id: str) -> bytes | None:
    """Load template file bytes into cache (read once, reuse many times)."""
    if theme_id in _template_bytes_cache:
        return _template_bytes_cache[theme_id]

    path = _resolve_template_path(theme_id)
    if not path:
        logger.warning(f"Template file not found for theme: {theme_id}")
        return None

    logger.info(f"Loading template into cache: {path} ({path.stat().st_size / 1024:.0f} KB)")
    data = path.read_bytes()
    _template_bytes_cache[theme_id] = data
    return data


def _open_template(theme_id: str) -> Presentation | None:
    """Open a fresh copy of the cached template (fast: reads from memory)."""
    data = _load_template_bytes(theme_id)
    if data is None:
        return None
    return Presentation(io.BytesIO(data))


# ── Layout Detection ────────────────────────────────────────

# Keywords to identify layout types (ordered by priority)
_TITLE_KEYWORDS = ["title slide", "title", "cover", "intro", "opening", "bìa", "tiêu đề"]
_CONTENT_KEYWORDS = ["content", "body", "text", "two content", "comparison",
                      "nội dung", "blank", "custom"]
_SECTION_KEYWORDS = ["section", "divider", "header", "phần"]
_ENDING_KEYWORDS = ["thank", "end", "closing", "kết", "cảm ơn"]


def _score_layout_type(layout_name: str, keywords: list[str]) -> int:
    """Score how well a layout name matches a list of keywords."""
    name = layout_name.lower().strip()
    for i, kw in enumerate(keywords):
        if kw in name:
            return len(keywords) - i  # Higher score for earlier keywords
    return 0


def _detect_layouts(prs: Presentation) -> dict:
    """
    Auto-detect the best layouts for: title, content, section, ending.
    Returns dict with layout indices.
    """
    layouts = list(prs.slide_layouts)
    if not layouts:
        return {"title": 0, "content": 0, "section": 0, "ending": 0}

    best = {"title": None, "content": None, "section": None, "ending": None}
    best_scores = {"title": 0, "content": 0, "section": 0, "ending": 0}

    for idx, layout in enumerate(layouts):
        name = layout.name

        # Score for each type
        for ltype, keywords in [
            ("title", _TITLE_KEYWORDS),
            ("content", _CONTENT_KEYWORDS),
            ("section", _SECTION_KEYWORDS),
            ("ending", _ENDING_KEYWORDS),
        ]:
            score = _score_layout_type(name, keywords)
            if score > best_scores[ltype]:
                best_scores[ltype] = score
                best[ltype] = idx

    # Fallback: if no match found, use heuristic indices
    # Most templates: index 0 = title, index 1 = content
    if best["title"] is None:
        best["title"] = 0
    if best["content"] is None:
        # Use the second layout, or the first one with most placeholders
        if len(layouts) > 1:
            best["content"] = 1
        else:
            best["content"] = 0
    if best["section"] is None:
        best["section"] = best["title"]  # Reuse title layout for sections
    if best["ending"] is None:
        best["ending"] = best["title"]  # Reuse title layout for ending

    layout_names = {k: layouts[v].name for k, v in best.items()}
    logger.info(f"Detected layouts: {layout_names}")

    return best


def _find_placeholders(slide):
    """
    Find title and body text areas in a slide.
    Uses a 3-tier strategy:
    1. Standard placeholders (idx 0 = title, idx 1 = body)
    2. Any placeholder with a text frame
    3. Any shape with a text frame (text boxes, etc.)
    """
    title_ph = None
    body_ph = None

    # --- Tier 1: Standard placeholder indices ---
    try:
        for shape in slide.placeholders:
            idx = shape.placeholder_format.idx
            if idx == 0:
                title_ph = shape
            elif idx == 1:
                body_ph = shape
            elif idx in (10, 11, 12, 13, 14, 15, 16) and body_ph is None:
                body_ph = shape
    except Exception:
        pass

    if title_ph and body_ph:
        return title_ph, body_ph

    # --- Tier 2: Any placeholder with text frame ---
    try:
        for shape in slide.placeholders:
            if not shape.has_text_frame:
                continue
            if title_ph is None:
                title_ph = shape
            elif body_ph is None and shape != title_ph:
                body_ph = shape
    except Exception:
        pass

    if title_ph and body_ph:
        return title_ph, body_ph

    # --- Tier 3: Any shape with a text frame (sorted by area, largest = body) ---
    text_shapes = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape != title_ph and shape != body_ph:
            area = shape.width * shape.height if shape.width and shape.height else 0
            text_shapes.append((area, shape))

    # Sort by area descending — largest text frame is likely body
    text_shapes.sort(key=lambda x: x[0], reverse=True)

    for area, shape in text_shapes:
        if title_ph is None:
            title_ph = shape
        elif body_ph is None and shape != title_ph:
            body_ph = shape
            break

    logger.debug(f"Found placeholders: title={title_ph is not None}, body={body_ph is not None}")
    return title_ph, body_ph


def _fill_text_frame(text_frame, text, max_font_size=None, is_title=False):
    """
    Fill a text frame with text.
    Caps font size to prevent overflow and enables auto-fit.
    """
    if not text_frame or not text:
        return

    # Default max font sizes
    if max_font_size is None:
        max_font_size = Pt(28) if is_title else Pt(14)

    # Try to preserve the formatting of the first paragraph/run
    existing_font = None
    if text_frame.paragraphs:
        for para in text_frame.paragraphs:
            for run in para.runs:
                try:
                    color_rgb = run.font.color.rgb if run.font.color.type is not None else None
                except Exception:
                    color_rgb = None
                existing_font = {
                    "name": run.font.name,
                    "size": run.font.size,
                    "bold": run.font.bold,
                    "italic": run.font.italic,
                    "color": color_rgb,
                }
                break
            if existing_font:
                break

    # Clear existing text
    text_frame.clear()
    text_frame.word_wrap = True

    # Enable auto-fit so text shrinks if needed
    try:
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass

    # Split text into paragraphs
    lines = text.split("\n")
    first_para = True

    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue

        para = text_frame.paragraphs[0] if first_para else text_frame.add_paragraph()
        first_para = False

        # Handle bullet points
        is_bullet = stripped.startswith(("- ", "• ", "* ", "– "))
        if is_bullet:
            stripped = stripped.lstrip("-•*– ").strip()
            stripped = f"•  {stripped}"

        run = para.add_run()
        run.text = stripped

        # Determine font size: use existing but cap at max
        font_size = max_font_size
        if existing_font and existing_font["size"]:
            font_size = min(existing_font["size"], max_font_size)

        # Apply formatting
        if existing_font and existing_font["name"]:
            run.font.name = existing_font["name"]
        run.font.size = font_size
        if is_title:
            run.font.bold = True
        elif existing_font and existing_font["bold"] is not None:
            run.font.bold = existing_font["bold"]
        if existing_font and existing_font["color"]:
            try:
                run.font.color.rgb = existing_font["color"]
            except Exception:
                pass

        # Spacing between paragraphs for body text
        if not is_title:
            para.space_before = Pt(3)
            para.space_after = Pt(3)


def _clear_all_slides(prs: Presentation):
    """Remove all existing slides from the presentation, keeping layouts intact."""
    sldIdLst = prs.slides._sldIdLst
    sldId_elements = list(sldIdLst)
    for sldId in sldId_elements:
        try:
            rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
            if rId is None:
                rId = sldId.rId
            prs.part.drop_rel(rId)
        except Exception as e:
            logger.warning(f"Could not drop relationship: {e}")
        sldIdLst.remove(sldId)
    logger.info(f"Cleared all slides from template ({len(sldId_elements)} removed)")


# ── Main API ────────────────────────────────────────────────

def build_from_template(
    theme_id: str,
    slides_data: list[dict],
    image_paths: dict | None = None,
) -> Presentation | None:
    """
    Build a presentation from a downloaded .pptx template.

    Strategy: KEEP the original template slides (each has unique design).
    - Replace text in existing slides with user content.
    - If user needs more slides → duplicate from template using their layouts.
    - If user needs fewer slides → remove extra template slides.

    Args:
        theme_id: Theme identifier (e.g. 'corporate_blue', 'neon_pop')
        slides_data: List of dicts with: slide_number, title, content, narration
        image_paths: Optional dict mapping slide_number -> image file path

    Returns:
        Presentation object, or None if template not found
    """
    prs = _open_template(theme_id)
    if prs is None:
        return None

    if not slides_data:
        return prs

    num_template_slides = len(prs.slides)
    num_user_slides = len(slides_data)
    logger.info(f"Template '{theme_id}': {num_template_slides} slides, user wants {num_user_slides}")

    # --- Step 1: If user needs MORE slides than template, add extras ---
    if num_user_slides > num_template_slides and num_template_slides > 0:
        # Add slides by reusing template slide layouts (cycling)
        for extra_i in range(num_user_slides - num_template_slides):
            # Cycle through existing template slides' layouts
            source_idx = (extra_i + 1) % num_template_slides  # skip cover (idx 0)
            if source_idx == 0 and num_template_slides > 1:
                source_idx = 1
            source_layout = prs.slides[source_idx].slide_layout
            prs.slides.add_slide(source_layout)
        logger.info(f"Added {num_user_slides - num_template_slides} extra slides")

    # --- Step 2: Replace text in each slide ---
    for i, slide_data in enumerate(slides_data):
        if i >= len(prs.slides):
            break

        slide = prs.slides[i]
        title_text = slide_data.get("title", "")
        content_text = slide_data.get("content", "")
        narration_text = slide_data.get("narration", "")
        is_title_slide = (i == 0)

        # Find text areas in this slide
        title_ph, body_ph = _find_placeholders(slide)

        # === Fill TITLE ===
        if title_ph and title_ph.has_text_frame:
            _fill_text_frame(title_ph.text_frame, title_text, is_title=True)
        elif title_text:
            # Fallback: create textbox
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.3),
                prs.slide_width - Inches(1.0), Inches(1.2)
            )
            _fill_text_frame(title_box.text_frame, title_text, is_title=True)

        # === Fill BODY ===
        if body_ph and body_ph.has_text_frame:
            if is_title_slide:
                short = content_text.split("\n")[0][:150] if content_text else ""
                _fill_text_frame(body_ph.text_frame, short, max_font_size=Pt(18))
            else:
                _fill_text_frame(body_ph.text_frame, content_text, is_title=False)
        elif content_text and not is_title_slide:
            # Fallback: create textbox
            body_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.6),
                prs.slide_width - Inches(1.0), prs.slide_height - Inches(2.2)
            )
            _fill_text_frame(body_box.text_frame, content_text, is_title=False)
        elif content_text and is_title_slide:
            short = content_text.split("\n")[0][:150]
            sub_box = slide.shapes.add_textbox(
                Inches(1.0), Inches(1.8),
                prs.slide_width - Inches(2.0), Inches(1.0)
            )
            _fill_text_frame(sub_box.text_frame, short, max_font_size=Pt(16))

        # Speaker notes
        if narration_text:
            try:
                notes_tf = slide.notes_slide.notes_text_frame
                notes_tf.text = narration_text
            except Exception:
                pass

    # --- Step 3: Remove extra template slides if user needs fewer ---
    while len(prs.slides) > num_user_slides:
        _remove_last_slide(prs)

    logger.info(f"Built presentation from template '{theme_id}': {len(prs.slides)} slides")
    return prs


def _remove_last_slide(prs: Presentation):
    """Remove the last slide from the presentation."""
    sldIdLst = prs.slides._sldIdLst
    if len(sldIdLst) == 0:
        return
    last = sldIdLst[-1]
    try:
        rId = last.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        if rId is None:
            rId = last.rId
        prs.part.drop_rel(rId)
    except Exception as e:
        logger.warning(f"Could not drop rel for last slide: {e}")
    sldIdLst.remove(last)


def preload_template(theme_id: str) -> bool:
    """Pre-load a template into memory cache. Returns True if successful."""
    data = _load_template_bytes(theme_id)
    return data is not None


def preload_all_templates():
    """Pre-load all templates into memory cache at startup."""
    loaded = 0
    for theme_id in AVAILABLE_THEMES:
        if preload_template(theme_id):
            loaded += 1
    logger.info(f"Pre-loaded {loaded}/{len(AVAILABLE_THEMES)} templates into cache")
    return loaded


def get_cached_count() -> int:
    """Return how many templates are currently cached."""
    return len(_template_bytes_cache)


def clear_cache():
    """Clear the template cache."""
    _template_bytes_cache.clear()
    logger.info("Template cache cleared")
