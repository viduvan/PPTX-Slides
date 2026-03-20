"""Inspect all template PPTX files to see current slide content."""
from pptx import Presentation
from pathlib import Path

templates_dir = Path(r'd:\AI_project\PPTX-Slides\templates')
for cat_dir in sorted(templates_dir.iterdir()):
    if not cat_dir.is_dir():
        continue
    print(f'\n=== {cat_dir.name} ===')
    for pptx_file in sorted(cat_dir.glob('*.pptx')):
        prs = Presentation(str(pptx_file))
        print(f'\n  {pptx_file.name} ({pptx_file.stat().st_size / 1024:.0f} KB)')
        print(f'    Slides: {len(prs.slides)}')
        for i, slide in enumerate(prs.slides):
            layout_name = slide.slide_layout.name if slide.slide_layout else 'N/A'
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    txt = shape.text_frame.text.strip()
                    if txt:
                        texts.append(txt[:100])
            print(f'    Slide {i+1} (layout: {layout_name}):')
            for t in texts[:5]:
                print(f'      "{t}"')
