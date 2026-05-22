"""
Exporter Agent Service — Builds PPTX/HTML files and generates thumbnails.
NO LLM calls — only file rendering and conversion.

HTML output uses the html-ppt design system:
  - 36 themes (CSS variable overrides)
  - 31 layout templates (single-page HTML)
  - 27 CSS animations + 20 Canvas FX
  - Keyboard runtime (← → T S O F navigation)

Developed by ChimSe (viduvan) - https://github.com/viduvan
Based on html-ppt-skill by lewis (MIT License)
"""
import json
import logging
import os
import re
import shutil
import subprocess
import tempfile
from pathlib import Path

from fastapi import FastAPI, HTTPException
from pydantic import BaseModel

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("exporter")

app = FastAPI(title="PPTX-Slides Exporter Agent", version="2.0.0")

SHARED_DATA_DIR = Path(os.getenv("SHARED_DATA_DIR", "/data"))
TEMPLATES_DIR = Path(os.getenv("TEMPLATES_DIR", "/app/templates"))
HTML_ASSETS_DIR = Path(os.getenv("HTML_ASSETS_DIR", "/app/html-assets"))
OUTPUT_DIR = SHARED_DATA_DIR / "output"
THUMBNAILS_DIR = SHARED_DATA_DIR / "thumbnails"

# Available themes and layouts (auto-detected at startup)
AVAILABLE_THEMES: list[str] = []
AVAILABLE_LAYOUTS: list[str] = []

# Layout type mapping: slide content → best layout template
LAYOUT_MAP = {
    "cover": "cover",
    "title": "cover",
    "toc": "toc",
    "table_of_contents": "toc",
    "bullets": "bullets",
    "bullet_points": "bullets",
    "two_column": "two-column",
    "three_column": "three-column",
    "comparison": "comparison",
    "pros_cons": "pros-cons",
    "timeline": "timeline",
    "roadmap": "roadmap",
    "process": "process-steps",
    "steps": "process-steps",
    "kpi": "kpi-grid",
    "metrics": "kpi-grid",
    "stats": "stat-highlight",
    "stat": "stat-highlight",
    "quote": "big-quote",
    "table": "table",
    "code": "code",
    "terminal": "terminal",
    "chart_bar": "chart-bar",
    "chart_line": "chart-line",
    "chart_pie": "chart-pie",
    "chart_radar": "chart-radar",
    "image": "image-hero",
    "image_grid": "image-grid",
    "gallery": "image-grid",
    "mindmap": "mindmap",
    "flow": "flow-diagram",
    "architecture": "arch-diagram",
    "gantt": "gantt",
    "checklist": "todo-checklist",
    "section": "section-divider",
    "divider": "section-divider",
    "cta": "cta",
    "thanks": "thanks",
    "end": "thanks",
    "diff": "diff",
}


@app.on_event("startup")
def _discover_assets():
    """Auto-detect available themes and layouts."""
    global AVAILABLE_THEMES, AVAILABLE_LAYOUTS
    themes_dir = HTML_ASSETS_DIR / "themes"
    layouts_dir = HTML_ASSETS_DIR / "layouts"

    if themes_dir.exists():
        AVAILABLE_THEMES = sorted([f.stem for f in themes_dir.glob("*.css")])
    if layouts_dir.exists():
        AVAILABLE_LAYOUTS = sorted([f.stem for f in layouts_dir.glob("*.html")])

    logger.info(f"Discovered {len(AVAILABLE_THEMES)} themes, {len(AVAILABLE_LAYOUTS)} layouts")


# ── Models ───────────────────────────────────────────────────

class BuildPptxRequest(BaseModel):
    job_id: str
    slides: list[dict]
    visual_specs: dict = {}


class BuildPptxResponse(BaseModel):
    job_id: str
    pptx_path: str
    slide_count: int


class BuildHtmlRequest(BaseModel):
    job_id: str
    slides: list[dict]
    visual_specs: dict = {}


class BuildHtmlResponse(BaseModel):
    job_id: str
    html_path: str
    assets_dir: str


class ThumbnailRequest(BaseModel):
    job_id: str
    pptx_path: str = ""
    html_path: str = ""


class ThumbnailResponse(BaseModel):
    job_id: str
    thumbnail_paths: list[str]
    count: int


# ── Endpoints ────────────────────────────────────────────────

@app.get("/health")
async def health():
    lo = shutil.which("libreoffice") or shutil.which("soffice")
    pdftoppm = shutil.which("pdftoppm")
    return {
        "status": "ok",
        "agent": "exporter",
        "version": "2.0.0",
        "libreoffice": bool(lo),
        "pdftoppm": bool(pdftoppm),
        "themes": len(AVAILABLE_THEMES),
        "layouts": len(AVAILABLE_LAYOUTS),
    }


@app.get("/themes")
async def list_themes():
    """List all available themes."""
    return {"themes": AVAILABLE_THEMES, "count": len(AVAILABLE_THEMES)}


@app.get("/layouts")
async def list_layouts():
    """List all available layout templates."""
    return {"layouts": AVAILABLE_LAYOUTS, "count": len(AVAILABLE_LAYOUTS)}


@app.post("/build-pptx", response_model=BuildPptxResponse)
async def build_pptx(req: BuildPptxRequest):
    """
    Build a PPTX file from slides + visual specs using python-pptx.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    job_dir = OUTPUT_DIR / req.job_id
    job_dir.mkdir(parents=True, exist_ok=True)

    theme = req.visual_specs.get("theme", "midnight")
    images = req.visual_specs.get("images", {})

    # Find template or create blank
    template_path = _find_pptx_template(theme)
    if template_path and template_path.exists():
        prs = Presentation(str(template_path))
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    for slide_data in req.slides:
        slide = prs.slides.add_slide(slide_layout)
        slide_num = slide_data.get("slide_number", 0)
        title = slide_data.get("title", "")
        content = slide_data.get("content", "")

        # Add title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.alignment = PP_ALIGN.LEFT

        # Add content bullets
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.8), Inches(7.5), Inches(5.2)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        lines = content.split("\n")
        for i, line in enumerate(lines):
            line = line.strip()
            if not line:
                continue
            if line.startswith("-"):
                line = line[1:].strip()

            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {line}"
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
            p.space_after = Pt(6)

        # Add image if available
        img_path = images.get(str(slide_num), "")
        if img_path and Path(img_path).exists():
            try:
                slide.shapes.add_picture(
                    img_path, Inches(8.5), Inches(1.8), Inches(4.3), Inches(4.3)
                )
            except Exception as e:
                logger.warning(f"Failed to add image for slide {slide_num}: {e}")

    # Save
    pptx_path = job_dir / f"presentation_{req.job_id}.pptx"
    prs.save(str(pptx_path))

    logger.info(f"[{req.job_id}] Built PPTX: {len(req.slides)} slides → {pptx_path}")
    return BuildPptxResponse(
        job_id=req.job_id,
        pptx_path=str(pptx_path),
        slide_count=len(req.slides),
    )


@app.post("/build-html", response_model=BuildHtmlResponse)
async def build_html(req: BuildHtmlRequest):
    """
    Build an HTML presentation deck using html-ppt design system.
    Uses real layout templates + theme CSS + animations + keyboard runtime.
    Supports 36 themes × 31 layouts × 27 animations.
    """
    job_dir = OUTPUT_DIR / req.job_id
    job_dir.mkdir(parents=True, exist_ok=True)

    theme = req.visual_specs.get("theme", "tokyo-night")
    layouts_spec = req.visual_specs.get("layouts", [])
    images = req.visual_specs.get("images", {})

    # Validate theme
    if theme not in AVAILABLE_THEMES:
        theme = "tokyo-night" if "tokyo-night" in AVAILABLE_THEMES else AVAILABLE_THEMES[0]

    # Copy assets to output dir for self-contained deck
    assets_dest = job_dir / "assets"
    if assets_dest.exists():
        shutil.rmtree(assets_dest)
    shutil.copytree(HTML_ASSETS_DIR, assets_dest, dirs_exist_ok=True)

    # Remove layout HTML files from assets (not needed in output)
    layouts_in_assets = assets_dest / "layouts"
    if layouts_in_assets.exists():
        shutil.rmtree(layouts_in_assets)

    total = len(req.slides)

    # Build slide sections
    sections_html = []
    for idx, slide_data in enumerate(req.slides):
        slide_num = slide_data.get("slide_number", idx + 1)
        title = slide_data.get("title", "")
        content = slide_data.get("content", "")
        narration = slide_data.get("narration", "")
        layout_type = _resolve_layout(slide_data, layouts_spec, idx, total)
        img_path = images.get(str(slide_num), "")

        section_html = _build_slide_section(
            title=title,
            content=content,
            narration=narration,
            layout_type=layout_type,
            slide_num=slide_num,
            total_slides=total,
            img_path=img_path,
            is_first=(idx == 0),
        )
        sections_html.append(section_html)

    # Build list of themes for T-key cycling (pick 6 representative ones)
    theme_cycle = [theme]
    cycle_candidates = ["tokyo-night", "aurora", "minimal-white", "dracula",
                        "catppuccin-mocha", "corporate-clean", "cyberpunk-neon"]
    for t in cycle_candidates:
        if t != theme and t in AVAILABLE_THEMES and len(theme_cycle) < 6:
            theme_cycle.append(t)

    deck_title = req.slides[0].get("title", "Presentation") if req.slides else "Presentation"

    html_content = f"""<!DOCTYPE html>
<html lang="en" data-theme="{theme}">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1.0">
<title>{_escape_html(deck_title)}</title>
<link rel="stylesheet" href="assets/fonts.css">
<link rel="stylesheet" href="assets/base.css">
<link rel="stylesheet" id="theme-link" href="assets/themes/{theme}.css">
<link rel="stylesheet" href="assets/animations/animations.css">
</head>
<body>
<div class="deck" data-themes="{','.join(theme_cycle)}" data-theme-base="assets/themes/">
{''.join(sections_html)}
</div>
<script src="assets/runtime.js"></script>
</body></html>"""

    html_path = job_dir / f"presentation_{req.job_id}.html"
    html_path.write_text(html_content, encoding="utf-8")

    logger.info(
        f"[{req.job_id}] Built HTML deck: {len(req.slides)} slides, "
        f"theme={theme}, → {html_path}"
    )
    return BuildHtmlResponse(
        job_id=req.job_id,
        html_path=str(html_path),
        assets_dir=str(assets_dest),
    )


@app.post("/thumbnails", response_model=ThumbnailResponse)
async def generate_thumbnails(req: ThumbnailRequest):
    """Generate PNG thumbnails from PPTX using LibreOffice + pdftoppm."""
    pptx_path = Path(req.pptx_path) if req.pptx_path else None
    if not pptx_path or not pptx_path.exists():
        raise HTTPException(status_code=404, detail="PPTX file not found")

    thumb_dir = THUMBNAILS_DIR / req.job_id
    thumb_dir.mkdir(parents=True, exist_ok=True)

    lo_cmd = shutil.which("libreoffice") or shutil.which("soffice") or "libreoffice"
    pdftoppm_cmd = shutil.which("pdftoppm") or "pdftoppm"

    try:
        with tempfile.TemporaryDirectory() as tmpdir:
            tmpdir = Path(tmpdir)

            # Step 1: PPTX → PDF
            result = subprocess.run(
                [lo_cmd, "--headless", "--convert-to", "pdf",
                 "--outdir", str(tmpdir), str(pptx_path)],
                capture_output=True, text=True, timeout=120,
            )
            if result.returncode != 0:
                logger.error(f"LibreOffice failed: {result.stderr}")
                return ThumbnailResponse(job_id=req.job_id, thumbnail_paths=[], count=0)

            pdf_path = tmpdir / f"{pptx_path.stem}.pdf"
            if not pdf_path.exists():
                return ThumbnailResponse(job_id=req.job_id, thumbnail_paths=[], count=0)

            # Step 2: PDF → PNG
            subprocess.run(
                [pdftoppm_cmd, "-png", "-r", "150", "-scale-to-x", "1280",
                 "-scale-to-y", "-1", str(pdf_path), str(tmpdir / "slide")],
                capture_output=True, text=True, timeout=120,
            )

            # Step 3: Move PNGs
            paths = []
            for i, png in enumerate(sorted(tmpdir.glob("slide-*.png")), 1):
                dest = thumb_dir / f"slide_{i}.png"
                shutil.move(str(png), str(dest))
                paths.append(str(dest))

            logger.info(f"[{req.job_id}] Generated {len(paths)} thumbnails")
            return ThumbnailResponse(
                job_id=req.job_id,
                thumbnail_paths=paths,
                count=len(paths),
            )
    except subprocess.TimeoutExpired:
        logger.error(f"Thumbnail generation timed out for {req.job_id}")
        return ThumbnailResponse(job_id=req.job_id, thumbnail_paths=[], count=0)
    except Exception as e:
        logger.error(f"Thumbnail error: {e}")
        return ThumbnailResponse(job_id=req.job_id, thumbnail_paths=[], count=0)


# ── Slide Building Helpers ───────────────────────────────────

def _build_slide_section(
    title: str, content: str, narration: str, layout_type: str,
    slide_num: int, total_slides: int, img_path: str, is_first: bool,
) -> str:
    """Build a single <section class='slide'> from slide data + layout type."""
    active = ' is-active' if is_first else ''
    anim = ' data-anim="fade-up"' if not is_first else ''

    bullets = [line.strip().lstrip("- ") for line in content.split("\n") if line.strip()]

    # Notes section (for presenter mode)
    notes_html = f'\n    <div class="notes">{_escape_html(narration)}</div>' if narration else ''

    # Footer with slide numbers
    footer = (
        f'<div class="deck-footer">'
        f'<span class="dim2">PPTX-Slides</span>'
        f'<span class="slide-number" data-current="{slide_num}" data-total="{total_slides}"></span>'
        f'</div>'
    )

    # --- Layout-specific HTML ---
    if layout_type == "cover":
        pills = ""
        if bullets:
            pills_items = " ".join(f'<span class="pill">{b}</span>' for b in bullets[:5])
            pills = f'<div class="row wrap mt-l">{pills_items}</div>'
        return f"""
  <section class="slide{active}" data-title="{_escape_html(title)}">
    <div class="deck-header"><span class="eyebrow">Presentation</span></div>
    <div class="anim-stagger-list">
      <h1 class="h1{' anim-fade-up' if not is_first else ''}">{_escape_html(title)}</h1>
      <p class="lede">{_escape_html(bullets[0]) if bullets else ''}</p>
      {pills}
    </div>
    {footer}{notes_html}
  </section>"""

    elif layout_type == "section-divider":
        subtitle = bullets[0] if bullets else ""
        return f"""
  <section class="slide{active}" data-title="{_escape_html(title)}" style="text-align:center">
    <h1 class="h1 anim-rise-in" style="font-size:96px">{_escape_html(title)}</h1>
    <p class="lede" style="margin:0 auto">{_escape_html(subtitle)}</p>
    {footer}{notes_html}
  </section>"""

    elif layout_type == "two-column":
        mid = len(bullets) // 2
        left = bullets[:mid] if mid > 0 else bullets[:1]
        right = bullets[mid:] if mid > 0 else bullets[1:]
        left_html = "\n".join(f"        <li>{_escape_html(b)}</li>" for b in left)
        right_html = "\n".join(f"        <li>{_escape_html(b)}</li>" for b in right)
        return f"""
  <section class="slide{active}" data-title="{_escape_html(title)}">
    <h2 class="h2"{anim}>{_escape_html(title)}</h2>
    <div class="grid g2 mt-l">
      <div class="card"><ul class="anim-stagger-list" style="list-style:none;padding:0">{left_html}</ul></div>
      <div class="card"><ul class="anim-stagger-list" style="list-style:none;padding:0">{right_html}</ul></div>
    </div>
    {footer}{notes_html}
  </section>"""

    elif layout_type == "comparison":
        mid = len(bullets) // 2
        left = bullets[:mid] if mid > 0 else bullets[:1]
        right = bullets[mid:] if mid > 0 else bullets[1:]
        left_html = "\n".join(f"        <li>{_escape_html(b)}</li>" for b in left)
        right_html = "\n".join(f"        <li>{_escape_html(b)}</li>" for b in right)
        return f"""
  <section class="slide{active}" data-title="{_escape_html(title)}">
    <h2 class="h2"{anim}>{_escape_html(title)}</h2>
    <div style="display:grid;grid-template-columns:1fr 60px 1fr;gap:20px;margin-top:30px;align-items:stretch">
      <div class="card" style="border-top:3px solid var(--bad);padding:30px">
        <h3>Before</h3><ul style="padding-left:20px;color:var(--text-2)">{left_html}</ul>
      </div>
      <div style="display:flex;align-items:center;justify-content:center;font-size:48px;color:var(--text-3)">→</div>
      <div class="card" style="border-top:3px solid var(--good);padding:30px">
        <h3>After</h3><ul style="padding-left:20px;color:var(--text-2)">{right_html}</ul>
      </div>
    </div>
    {footer}{notes_html}
  </section>"""

    elif layout_type in ("kpi-grid", "stat-highlight"):
        cards_html = []
        for b in bullets[:6]:
            parts = b.split(":", 1) if ":" in b else (b, "")
            label = parts[0].strip()
            value = parts[1].strip() if len(parts) > 1 else ""
            cards_html.append(
                f'<div class="card"><p class="eyebrow">{_escape_html(label)}</p>'
                f'<div style="font-size:48px;font-weight:800">{_escape_html(value or label)}</div></div>'
            )
        cols = min(len(cards_html), 4)
        return f"""
  <section class="slide{active}" data-title="{_escape_html(title)}">
    <h2 class="h2"{anim}>{_escape_html(title)}</h2>
    <div class="grid g{cols} mt-l anim-stagger-list">
      {''.join(cards_html)}
    </div>
    {footer}{notes_html}
  </section>"""

    elif layout_type == "big-quote":
        quote_text = bullets[0] if bullets else title
        author = bullets[1] if len(bullets) > 1 else ""
        return f"""
  <section class="slide{active}" data-title="Quote" style="text-align:center">
    <blockquote style="font-family:var(--font-serif);font-size:42px;line-height:1.5;max-width:50ch;margin:0 auto;font-weight:300">
      "{_escape_html(quote_text)}"
    </blockquote>
    <p class="dim mt-l" style="font-size:18px">— {_escape_html(author)}</p>
    {footer}{notes_html}
  </section>"""

    elif layout_type == "thanks":
        subtitle = bullets[0] if bullets else "Thank you for your attention"
        return f"""
  <section class="slide{active}" data-title="Thanks" style="text-align:center">
    <h1 class="h1 gradient-text anim-zoom-pop" style="font-size:96px">{_escape_html(title)}</h1>
    <p class="lede" style="margin:18px auto 0">{_escape_html(subtitle)}</p>
    {footer}{notes_html}
  </section>"""

    elif layout_type == "image-hero" and img_path:
        return f"""
  <section class="slide{active}" data-title="{_escape_html(title)}" style="padding:0;position:relative">
    <div style="position:absolute;inset:0;background:url('{img_path}') center/cover;filter:brightness(0.4)"></div>
    <div style="position:relative;z-index:1;padding:72px 96px">
      <h1 class="h1" style="color:white">{_escape_html(title)}</h1>
      <p style="color:rgba(255,255,255,0.8);font-size:20px">{_escape_html(bullets[0]) if bullets else ''}</p>
    </div>
    {footer}{notes_html}
  </section>"""

    # --- Default: Bullets layout (most common) ---
    else:
        bullets_html = "\n".join(
            f'    <li class="card card-accent"><h4>{_escape_html(b)}</h4></li>'
            for b in bullets
        )
        return f"""
  <section class="slide{active}" data-title="{_escape_html(title)}">
    <h2 class="h2"{anim}>{_escape_html(title)}</h2>
    <ul class="grid g1 anim-stagger-list" style="list-style:none;padding:0;margin:20px 0 0;gap:14px">
{bullets_html}
    </ul>
    {footer}{notes_html}
  </section>"""


def _resolve_layout(slide_data: dict, layouts_spec: list, idx: int, total: int) -> str:
    """Determine which layout to use for a slide."""
    # 1. Explicit layout from visual_specs
    if idx < len(layouts_spec):
        layout_name = layouts_spec[idx]
        if layout_name in LAYOUT_MAP:
            return LAYOUT_MAP[layout_name]
        if layout_name in AVAILABLE_LAYOUTS:
            return layout_name

    # 2. Layout hint from slide data
    layout_hint = slide_data.get("layout", "").lower().replace(" ", "_")
    if layout_hint in LAYOUT_MAP:
        return LAYOUT_MAP[layout_hint]

    # 3. Auto-detect from position
    if idx == 0:
        return "cover"
    if idx == total - 1:
        title_lower = slide_data.get("title", "").lower()
        if any(w in title_lower for w in ["thank", "q&a", "end", "conclusion"]):
            return "thanks"

    # 4. Default to bullets
    return "bullets"


# ── General Helpers ──────────────────────────────────────────

def _escape_html(text: str) -> str:
    """Escape HTML special characters."""
    return (text
            .replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
            .replace('"', "&quot;"))


def _find_pptx_template(theme: str) -> Path | None:
    """Find PPTX template file for a theme."""
    for category in ["business", "creative", "education", "tech"]:
        path = TEMPLATES_DIR / category / f"{theme}.pptx"
        if path.exists():
            return path
    # Fallback: any .pptx in templates dir
    for pptx in TEMPLATES_DIR.rglob("*.pptx"):
        return pptx
    return None
