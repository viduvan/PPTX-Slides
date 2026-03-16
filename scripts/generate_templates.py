"""
Generate Template Files — Creates 30 .pptx template files in templates/ directory.
Run from project root:  python scripts/generate_templates.py
"""
import importlib.util
import sys
from pathlib import Path

# Project root
ROOT = Path(__file__).resolve().parent.parent
TEMPLATES_DIR = ROOT / "templates"

# Load template_builder directly (avoids heavy api.__init__ chain)
spec = importlib.util.spec_from_file_location(
    "template_builder",
    ROOT / "api" / "services" / "template_builder.py",
)
tb = importlib.util.module_from_spec(spec)
spec.loader.exec_module(tb)


from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

SLIDE_WIDTH = tb.SLIDE_WIDTH
SLIDE_HEIGHT = tb.SLIDE_HEIGHT
FONT_TITLE = tb.FONT_TITLE
FONT_BODY = tb.FONT_BODY
THEMES = tb.THEMES
THEME_REGISTRY = tb.THEME_REGISTRY
THEME_CATEGORIES = tb.THEME_CATEGORIES
AVAILABLE_THEMES = tb.AVAILABLE_THEMES


def _grad(slide, colors):
    bg = slide.background.fill
    bg.gradient(); bg.gradient_angle = 315
    bg.gradient_stops[0].position = 0.0
    bg.gradient_stops[0].color.rgb = colors["bg_dark"]
    bg.gradient_stops[1].position = 1.0
    bg.gradient_stops[1].color.rgb = colors["bg_gradient"]


def _bar(slide, colors, y=None):
    if y is None:
        y = SLIDE_HEIGHT - Inches(0.38)
    bar = slide.shapes.add_shape(1, Inches(0.5), y, SLIDE_WIDTH - Inches(1.0), Inches(0.08))
    bar.line.fill.background()
    f = bar.fill; f.gradient(); f.gradient_angle = 0
    f.gradient_stops[0].position = 0.0; f.gradient_stops[0].color.rgb = colors["accent"]
    f.gradient_stops[1].position = 1.0; f.gradient_stops[1].color.rgb = colors["accent_light"]


def build(theme_id):
    c = THEMES[theme_id]; reg = THEME_REGISTRY[theme_id]
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH; prs.slide_height = SLIDE_HEIGHT

    # Slide 1 — Title
    s = prs.slides.add_slide(prs.slide_layouts[6]); _grad(s, c)
    d = s.shapes.add_shape(1, Inches(.3), Inches(.3), Inches(.5), Pt(4))
    d.line.fill.background(); d.fill.solid(); d.fill.fore_color.rgb = c["accent"]
    b = s.shapes.add_textbox(Inches(1), Inches(2), Inches(11.3), Inches(2))
    tf = b.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "Presentation Title"
    r.font.name = FONT_TITLE; r.font.size = Pt(40); r.font.color.rgb = c["title"]; r.font.bold = True
    b2 = s.shapes.add_textbox(Inches(2), Inches(4.2), Inches(9.3), Inches(1))
    tf2 = b2.text_frame; tf2.word_wrap = True; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = f"Template: {reg['label']}"
    r2.font.name = FONT_BODY; r2.font.size = Pt(20); r2.font.color.rgb = c["subtitle"]
    _bar(s, c, Inches(5.5))

    # Slide 2 — Content
    s = prs.slides.add_slide(prs.slide_layouts[6]); _grad(s, c)
    b = s.shapes.add_textbox(Inches(.8), Inches(.5), Inches(11.5), Inches(1))
    tf = b.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "Content Slide Title"
    r.font.name = FONT_TITLE; r.font.size = Pt(28); r.font.color.rgb = c["title"]; r.font.bold = True
    ul = s.shapes.add_shape(1, Inches(.8), Inches(1.45), Inches(2), Pt(3))
    ul.line.fill.background(); ul.fill.solid(); ul.fill.fore_color.rgb = c["accent"]
    bb = s.shapes.add_textbox(Inches(.8), Inches(1.7), Inches(11.5), Inches(5))
    tfb = bb.text_frame; tfb.word_wrap = True; tfb.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, txt in enumerate(["Sample point one", "Key idea two", "Styled bullets", "Auto-fit sizing", "Gradient backgrounds"]):
        pa = tfb.paragraphs[0] if i == 0 else tfb.add_paragraph()
        rb = pa.add_run(); rb.text = f"  •  {txt}"
        rb.font.name = FONT_BODY; rb.font.size = Pt(16); rb.font.color.rgb = c["body"]
        pa.space_before = Pt(4); pa.space_after = Pt(4)
    _bar(s, c)

    # Slide 3 — Ending
    s = prs.slides.add_slide(prs.slide_layouts[6]); _grad(s, c)
    b = s.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(2))
    tf = b.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = "Thank You"
    r.font.name = FONT_TITLE; r.font.size = Pt(36); r.font.color.rgb = c["title"]; r.font.bold = True
    b2 = s.shapes.add_textbox(Inches(2), Inches(4.5), Inches(9.3), Inches(.8))
    tf2 = b2.text_frame; tf2.word_wrap = True; p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = "Questions? Contact us"
    r2.font.name = FONT_BODY; r2.font.size = Pt(18); r2.font.color.rgb = c["subtitle"]
    _bar(s, c, Inches(5.8))
    return prs


def main():
    print("=" * 60)
    print("  PPTX-Slides — Template Generator")
    print("=" * 60)
    n = 0
    for cat_id in sorted(THEME_CATEGORIES, key=lambda k: THEME_CATEGORIES[k]["order"]):
        ci = THEME_CATEGORIES[cat_id]
        d = TEMPLATES_DIR / cat_id; d.mkdir(parents=True, exist_ok=True)
        print(f"\n📁 {ci['emoji']} {ci['label']} ({cat_id}/)")
        for tid in AVAILABLE_THEMES:
            reg = THEME_REGISTRY.get(tid)
            if not reg or reg["category"] != cat_id:
                continue
            prs = build(tid)
            out = d / f"{tid}.pptx"; prs.save(str(out))
            kb = out.stat().st_size / 1024
            print(f"   ✅ {reg['emoji']} {reg['label']:20s} → {tid}.pptx ({kb:.0f} KB)")
            n += 1
    print(f"\n{'=' * 60}")
    print(f"  ✨ Generated {n} templates in {len(THEME_CATEGORIES)} categories!")
    print(f"  📂 Location: {TEMPLATES_DIR}")
    print(f"{'=' * 60}")


if __name__ == "__main__":
    main()
